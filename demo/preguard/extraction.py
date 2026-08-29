"""Raw bytes -> plain text, per supported document type."""
import io

import docx
import pptx
import pypdf


class UnsupportedFileType(ValueError):
    pass


def extract_text(raw: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(raw)
    if ext == "pptx":
        return _extract_pptx(raw)
    if ext == "docx":
        return _extract_docx(raw)
    if ext in ("txt", "csv", "md"):
        return raw.decode("utf-8", errors="replace")

    raise UnsupportedFileType(f"No extractor registered for .{ext}")


def _extract_pdf(raw: bytes) -> str:
    reader = pypdf.PdfReader(io.BytesIO(raw))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_pptx(raw: bytes) -> str:
    prs = pptx.Presentation(io.BytesIO(raw))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def _extract_docx(raw: bytes) -> str:
    d = docx.Document(io.BytesIO(raw))
    return "\n".join(p.text for p in d.paragraphs)
